import os
import convertapi
from pptx import Presentation

convertapi.api_credentials = os.getenv("CONVERTAPI_KEY")


def extract_text_slidewise(ppt_path):
    prs = Presentation(ppt_path)
    slides = []

    for index, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    if p.text.strip():
                        texts.append(p.text.strip())

        slides.append({
            "slide_number": index,
            "text": texts
        })

    return slides


def convert_ppt_to_images(ppt_path, output_dir):
    """
    IMPORTANT:
    - Do NOT trust filenames
    - Use creation time order
    - This guarantees correct slide mapping
    """
    result = convertapi.convert("jpg", {"File": ppt_path})
    result.save_files(output_dir)

    images = [
        os.path.join(output_dir, f)
        for f in os.listdir(output_dir)
        if f.lower().endswith(".jpg")
    ]

    # 🔥 CORE FIX: order by creation time (slide order)
    images.sort(key=os.path.getctime)

    return images
